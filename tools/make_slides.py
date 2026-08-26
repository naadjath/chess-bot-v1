"""Genere la presentation de soutenance (.pptx) et le guide oral (.md).

Contrairement a une premiere version texte-seule, cette presentation integre :
  - les VRAIES courbes d'entrainement et le VRAI tableau d'Elo, comme images
    (generees par tools/make_slide_assets.py a partir des donnees du projet) ;
  - des schemas dessines nativement (pipeline, encodage, architecture,
    masquage des coups illegaux) plutot que des listes a puces ;
  - un texte de fond precis (chiffres reels, mecanismes concrets), et un
    script oral detaille dans les notes de chaque diapo.

Prerequis : lancer d'abord `python tools/make_slide_assets.py`.

    python tools/make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "build" / "slide_assets"
#: Le pptx est un outil de travail (a presenter, modifiable) : il vit a la
#: racine du projet, pas dans LIVRABLES. La version PDF, elle, est le vrai
#: livrable et est copiee dans LIVRABLES par tools/build_livrables.py.
PPTX_OUT = ROOT / "Soutenance_Chess_Bot_v1.pptx"
#: Le guide oral est un outil de preparation personnel, jamais un livrable :
#: il reste a la racine du projet, jamais copie dans LIVRABLES.
ORAL_OUT = ROOT / "GUIDE-ORAL.md"

INK = RGBColor(0x1C, 0x1A, 0x1B)
SOFT = RGBColor(0x5D, 0x55, 0x59)
ROSE = RGBColor(0xB8, 0x53, 0x6E)
ROSE_LIGHT = RGBColor(0xE7, 0xC6, 0xCE)
LIGHT = RGBColor(0xF1, 0xEC, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRID = RGBColor(0xDD, 0xD6, 0xD9)

SLIDE_W = Emu(12192000)   # 16:9
SLIDE_H = Emu(6858000)
FONT = "Segoe UI"

MARGIN = Emu(700000)
CONTENT_TOP = Emu(1750000)

# Trace (titre, texte oral) de chaque diapo au fur et a mesure de la
# construction, pour regenerer le guide oral SANS jamais relire le pptx (donc
# sans jamais se tromper de titre ni oublier une diapo).
_ORAL_SCRIPT: list[tuple[str, str]] = []


def _note(slide, title, oral):
    """Enregistre le texte oral dans les notes PowerPoint ET dans le guide."""
    slide.notes_slide.notes_text_frame.text = "A DIRE :\n\n" + oral
    _ORAL_SCRIPT.append((title, oral))


# ---------------------------------------------------------------------------
# Primitives bas niveau
# ---------------------------------------------------------------------------

def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _run(paragraph, text, size, color, bold=False, italic=False, font=FONT):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _header(slide, title, subtitle=None, kicker=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Emu(520000), Emu(90000), Emu(760000))
    bar.fill.solid(); bar.fill.fore_color.rgb = ROSE
    bar.line.fill.background(); bar.shadow.inherit = False

    tf = _textbox(slide, Emu(920000), Emu(470000), Emu(10500000), Emu(1150000))
    if kicker:
        pk = tf.paragraphs[0]
        _run(pk, kicker.upper(), 12, ROSE, bold=True)
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    _run(p, title, 30, INK, bold=True)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        _run(p2, subtitle, 15, SOFT)


def _footer(slide, tag):
    tf = _textbox(slide, MARGIN, Emu(6480000), Emu(10800000), Emu(300000))
    _run(tf.paragraphs[0], f"Chess Bot v1  ·  {tag}", 9.5, RGBColor(0xA5, 0x9D, 0xA1))


def _box(slide, left, top, width, height, text, fill, text_color=INK, size=13,
          bold=False, align=PP_ALIGN.CENTER, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_color=None):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color; sh.line.width = Pt(1.2)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(80000)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _run(p, line, size, text_color, bold=bold)
    return sh


def _arrow(slide, x1, y1, x2, y2, color=ROSE, weight=2.4):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    _force_arrowhead(conn)
    return conn


def _force_arrowhead(conn):
    """python-pptx a un support limite des pointes de fleche : on force en XML."""
    ln = conn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)


def _caption(slide, left, top, width, text, size=11.5, color=SOFT, align=PP_ALIGN.CENTER, italic=True):
    tf = _textbox(slide, left, top, width, Emu(500000))
    _run(tf.paragraphs[0], text, size, color, italic=italic)
    tf.paragraphs[0].alignment = align


def _picture_centered(slide, path, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    pw, ph = Emu(int(w * ratio)), Emu(int(h * ratio))
    left = Emu(int((SLIDE_W - pw) / 2))
    slide.shapes.add_picture(str(path), left, top, width=pw, height=ph)


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    return slide


# ---------------------------------------------------------------------------
# Diapos de contenu textuel (avec puces detaillees, pas des mots-cles vagues)
# ---------------------------------------------------------------------------

def slide_bullets(prs, title, subtitle, bullets, oral, kicker=None, tag=""):
    slide = _new_slide(prs)
    _header(slide, title, subtitle, kicker)
    tf = _textbox(slide, Emu(920000), CONTENT_TOP, Emu(10400000), Emu(4500000))
    for i, (head, body) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        dot = p.add_run(); dot.text = "▸  "; dot.font.size = Pt(16); dot.font.color.rgb = ROSE; dot.font.bold = True; dot.font.name = FONT
        h = p.add_run(); h.text = head; h.font.size = Pt(16); h.font.bold = True; h.font.color.rgb = INK; h.font.name = FONT
        if body:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(14)
            b = p2.add_run(); b.text = "     " + body; b.font.size = Pt(13.5); b.font.color.rgb = SOFT; b.font.name = FONT
    _footer(slide, tag)
    _note(slide, title, oral)
    return slide


def slide_cover(prs, title, subtitle, lines, oral):
    slide = _new_slide(prs)
    _bg(slide, INK)
    tf = _textbox(slide, Emu(900000), Emu(2450000), Emu(10400000), Emu(2500000), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _run(p, title, 50, WHITE, bold=True)
    if subtitle:
        p2 = tf.add_paragraph(); p2.space_before = Pt(12)
        _run(p2, subtitle, 21, ROSE_LIGHT, bold=True)
    for line in lines:
        pl = tf.add_paragraph(); pl.space_before = Pt(8)
        _run(pl, line, 14, RGBColor(0xC9, 0xC2, 0xC5))
    _note(slide, title, oral)
    return slide


def slide_image(prs, title, subtitle, image_path, caption, oral, kicker=None, tag="", max_h=Emu(4600000)):
    slide = _new_slide(prs)
    _header(slide, title, subtitle, kicker)
    _picture_centered(slide, image_path, CONTENT_TOP, Emu(10700000), max_h)
    if caption:
        _caption(slide, MARGIN, Emu(6180000), Emu(10800000), caption)
    _footer(slide, tag)
    _note(slide, title, oral)
    return slide


# ---------------------------------------------------------------------------
# Schemas natifs (pptx) : pipeline, entonnoir, encodage, architecture, masquage
# ---------------------------------------------------------------------------

def diagram_pipeline(prs):
    slide = _new_slide(prs)
    _header(slide, "Vue d'ensemble", "Quatre etapes, de la partie brute au coup joue", "Comment ca marche")

    steps = [
        ("1", "DONNEES", "1M de positions\nissues de Lichess"),
        ("2", "ENCODAGE", "chaque position\n-> 68 nombres"),
        ("3", "MODELE", "le Transformer\nnote 1968 coups"),
        ("4", "EVALUATION", "matchs juges\n-> Elo"),
    ]
    n = len(steps)
    box_w, box_h = Emu(2380000), Emu(1750000)
    gap = Emu(280000)
    total_w = box_w * n + gap * (n - 1)
    left0 = Emu(int((SLIDE_W - total_w) / 2))
    top = Emu(2650000)

    for i, (num, label, desc) in enumerate(steps):
        x = Emu(int(left0 + i * (box_w + gap)))
        fill = ROSE if i in (1, 3) else INK
        _box(slide, x, top, box_w, Emu(560000), label, fill, WHITE, size=15, bold=True)
        _box(slide, x, Emu(top + 620000), box_w, Emu(1130000), desc, LIGHT, INK, size=12.5)
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x + box_w/2 - 220000)), Emu(top - 420000), Emu(440000), Emu(440000))
        num_box.fill.solid(); num_box.fill.fore_color.rgb = WHITE
        num_box.line.color.rgb = ROSE; num_box.line.width = Pt(2)
        num_box.shadow.inherit = False
        ntf = num_box.text_frame; ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _run(ntf.paragraphs[0], num, 18, ROSE, bold=True); ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if i < n - 1:
            xa = Emu(int(x + box_w + gap * 0.15))
            xb = Emu(int(x + box_w + gap * 0.85))
            _arrow(slide, xa, Emu(top + 900000), xb, Emu(top + 900000))

    _caption(slide, MARGIN, Emu(5150000), Emu(10800000),
             "L'entrainement est du clonage comportemental : le reseau apprend a imiter le coup joue par un joueur fort dans chaque position.")
    _footer(slide, "Architecture")
    _note(slide, "Vue d'ensemble", (
        "Voici les quatre etapes de la chaine complete, du fichier de parties brutes jusqu'au coup "
        "joue par le bot. Premiere etape : on constitue un million de positions a partir de vraies "
        "parties de joueurs forts. Deuxieme etape : chaque position est transformee en 68 nombres, "
        "le seul langage que comprend un reseau de neurones. Troisieme etape, le coeur du projet : "
        "le Transformer regarde ces 68 nombres et attribue un score a chacun des 1968 coups "
        "possibles. Quatrieme etape : on evalue le resultat en faisant vraiment jouer le bot, et on "
        "en tire un score Elo. On va detailler chacune de ces quatre etapes sur les diapos suivantes."
    ))
    return slide


def diagram_funnel(prs):
    slide = _new_slide(prs)
    _header(slide, "Etape 1 : les donnees", "Filtrer pour apprendre d'un bon joueur, pas de n'importe qui", "Donnees")

    stages = [
        ("186 057 parties lues", "archive Lichess de janvier 2024, lue en flux (aucune decompression sur le disque)", WHITE, INK, GRID),
        ("Filtres appliques", "Elo des deux joueurs >= 2000  ·  cadence >= 180 s (exclut le bullet)  ·  partie terminee  ·  8 premiers demi-coups ignores (ouverture memorisee)", LIGHT, INK, None),
        ("14 456 parties gardees", "soit 7,8 % des parties lues", WHITE, INK, GRID),
        ("1 000 012 positions", "980 440 pour l'entrainement · 19 572 pour la validation (separees PAR PARTIE, jamais par position, pour eviter la fuite de donnees)", ROSE, WHITE, None),
    ]
    top = 1900000
    widths = [9800000, 9200000, 8600000, 10800000]
    heights = [560000, 780000, 560000, 780000]
    for i, ((head, body, fill, color, line), w, h) in enumerate(zip(stages, widths, heights)):
        bw, bh = Emu(w), Emu(h)
        left = Emu(int((SLIDE_W - bw) / 2))
        sh = _box(slide, left, Emu(top), bw, bh, "", fill, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_color=line)
        tf = sh.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, head, 15, color, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(3)
        _run(p2, body, 10.5, color if fill == ROSE else SOFT)
        top += h + 90000
        if i < len(stages) - 1:
            ay = top - 90000
            _arrow(slide, Emu(int(SLIDE_W/2)), Emu(ay - 60000), Emu(int(SLIDE_W/2)), Emu(ay + 30000))
    _footer(slide, "Donnees")
    _note(slide, "Etape 1 : les donnees", (
        "On part de l'archive publique de Lichess de janvier 2024, lue directement en flux depuis "
        "internet : on ne telecharge jamais le fichier complet sur le disque, on s'arrete des qu'on a "
        "assez de positions. Sur 186 057 parties lues, on n'en garde que 14 456, moins de 8 pourcent. "
        "Pourquoi si peu ? On filtre sur l'Elo des deux joueurs, au moins 2000, sur la cadence, pour "
        "exclure le bullet ou l'on joue au reflexe sans reflechir, et sur les huit premiers coups de "
        "chaque partie, qui relevent surtout de la memorisation d'ouvertures plutot que de la reflexion. "
        "Un dernier point technique important : on separe l'entrainement et la validation PAR PARTIE, "
        "jamais par position isolee. Deux positions issues de la meme partie se ressemblent enormement ; "
        "les melanger aurait fausse notre mesure en laissant le modele etre evalue sur des situations "
        "presque deja vues."
    ))
    return slide


def diagram_encoding(prs):
    slide = _new_slide(prs)
    _header(slide, "Etape 2 : encoder une position", "Une position d'echecs devient une sequence de 68 nombres", "Encodage")

    # Grille 8x8 a gauche (miniature)
    cell = 230000
    gx, gy = 950000, 2050000
    for r in range(8):
        for c in range(8):
            x = Emu(int(gx + c * cell)); y = Emu(int(gy + r * cell))
            color = LIGHT if (r + c) % 2 == 0 else SOFT
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(cell), Emu(cell))
            sq.fill.solid(); sq.fill.fore_color.rgb = color
            sq.line.fill.background(); sq.shadow.inherit = False
    _caption(slide, Emu(gx), Emu(gy + 8*cell + 60000), Emu(8*cell), "L'echiquier : 64 cases", size=11)

    _arrow(slide, Emu(gx + 8*cell + 150000), Emu(gy + 4*cell), Emu(gx + 8*cell + 750000), Emu(gy + 4*cell))

    # Bloc des 68 tokens a droite
    tx = gx + 8*cell + 900000
    rows = [
        ("tokens 0 - 63", "contenu de chaque case (0 = vide, 1-6 = pieces alliees, 7-12 = pieces adverses)"),
        ("token 64", "le trait (qui doit jouer)"),
        ("tokens 65 - 66", "droits de roque"),
        ("token 67", "case de prise en passant"),
    ]
    ty = gy
    for head, body in rows:
        bw, bh = 4700000, 560000
        sh = _box(slide, Emu(tx), Emu(ty), Emu(bw), Emu(bh), "", LIGHT, INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = sh.text_frame
        p = tf.paragraphs[0]; _run(p, head, 13, ROSE, bold=True)
        p2 = tf.add_paragraph(); p2.space_before = Pt(2); _run(p2, body, 10.5, SOFT)
        ty = ty + bh + 90000

    _box(slide, Emu(tx), Emu(ty), Emu(4700000), Emu(500000), "= 68 tokens en entree du reseau", ROSE, WHITE, size=13, bold=True)

    _caption(slide, MARGIN, Emu(5980000), Emu(10800000),
             "Astuce cle : l'echiquier est toujours retourne du point de vue du joueur au trait. Le modele n'apprend ainsi qu'un seul point de vue, deux fois plus vite.")
    _footer(slide, "Encodage")
    _note(slide, "Etape 2 : encoder une position", (
        "Un reseau de neurones ne comprend que des nombres, jamais un plateau. On transforme donc "
        "chaque position en 68 tokens : les 64 premiers decrivent le contenu de chaque case, zero pour "
        "une case vide, un a six pour les pieces du joueur au trait, sept a douze pour les pieces "
        "adverses. Les quatre derniers tokens portent le trait, les droits de roque et la prise en "
        "passant. Le detail le plus important a expliquer au jury : on retourne systematiquement "
        "l'echiquier pour que le modele voie toujours la position du point de vue du joueur qui doit "
        "jouer, comme s'il jouait toujours les blancs. Sans cette astuce, il faudrait apprendre deux "
        "fois la meme chose, une fois pour les blancs et une fois pour les noirs. Avec elle, chaque "
        "partie fournit deux fois plus d'exemples utiles."
    ))
    return slide


def diagram_architecture(prs):
    slide = _new_slide(prs)
    _header(slide, "Etape 3 : le modele", "Un Transformer encodeur de 6,86 millions de parametres", "Architecture")

    cx = int(SLIDE_W/2)
    top = 1850000

    _box(slide, Emu(cx - 1600000), Emu(top), Emu(3200000), Emu(480000), "Entree : 68 tokens", INK, WHITE, size=13, bold=True)
    y = top + 560000
    _arrow(slide, Emu(cx), Emu(top + 480000), Emu(cx), Emu(y + 20000))

    # Pile de 8 couches d'attention (dessinees comme 3 blocs visibles + indication x8)
    layer_h = 310000
    layer_w = 5200000
    y2 = y + 100000
    for i in range(3):
        yy = y2 + i * (layer_h + 70000)
        _box(slide, Emu(cx - layer_w//2), Emu(yy), Emu(layer_w), Emu(layer_h),
             "Couche d'attention" if i < 2 else "…",
             ROSE_LIGHT if i < 2 else WHITE, INK, size=11.5)
    stack_bottom = y2 + 3*(layer_h+70000)
    _caption(slide, Emu(cx - layer_w//2), Emu(stack_bottom + 30000), Emu(layer_w),
             "x 8 couches  ·  8 tetes d'attention par couche  ·  dimension 256", size=11.5, italic=False, color=SOFT)

    arrow_top = stack_bottom + 340000
    y3 = arrow_top + 340000
    _arrow(slide, Emu(cx), Emu(arrow_top), Emu(cx), Emu(y3))
    _box(slide, Emu(cx - 1600000), Emu(y3), Emu(3200000), Emu(480000), "Sortie : 1 score par coup (1968)", ROSE, WHITE, size=13, bold=True)

    # Colonne d'infos a droite
    info_x = cx + 3200000
    info = [
        ("6 858 416", "parametres au total"),
        ("Pre-normalisation", "stabilise l'entrainement d'un reseau profond"),
        ("Attention", "relie directement 2 cases eloignees - utile pour une tour ou un fou"),
    ]
    iy = 1900000
    for head, body in info:
        _box(slide, Emu(info_x), Emu(iy), Emu(2600000), Emu(950000), head + "\n" + body, LIGHT, INK, size=11.5, bold=False)
        iy = iy + 1050000

    _footer(slide, "Architecture")
    _note(slide, "Etape 3 : le modele", (
        "Voici le coeur du projet. Les 68 tokens entrent dans un Transformer encodeur, huit couches "
        "empilees, chacune avec huit tetes d'attention. Pourquoi un Transformer plutot qu'un reseau "
        "plus classique comme un CNN ? Parce que son mecanisme d'attention relie directement n'importe "
        "quelle case a n'importe quelle autre case, en une seule etape. C'est exactement le "
        "comportement d'une tour ou d'un fou, qui agissent a distance, alors qu'un reseau convolutif "
        "classique ne regarde d'abord que les cases voisines. En sortie, le modele produit un score "
        "pour chacun des 1968 coups possibles : c'est un probleme de classification, comme reconnaitre "
        "une image, sauf qu'ici on classe une position parmi 1968 categories de coup. Au total, "
        "6 858 416 parametres, ce qui reste tres modeste face aux modeles de DeepMind."
    ))
    return slide


def diagram_masking(prs):
    slide = _new_slide(prs)
    _header(slide, "Choisir un coup legal", "Le bot ne peut pas tricher : c'est impossible par construction", "Architecture")

    boxes = [
        ("Le reseau note\nles 1968 coups", LIGHT, INK),
        ("On repere les coups\nILLEGAUX dans\nla position actuelle", ROSE_LIGHT, INK),
        ("Leur score devient\n-l'infini", INK, WHITE),
        ("On garde le\nmeilleur score\nrestant", ROSE, WHITE),
    ]
    n = len(boxes)
    bw, bh = 2500000, 1500000
    gap = 280000
    total = bw * n + gap * (n - 1)
    left0 = int((SLIDE_W - total) / 2)
    top = 2350000
    for i, (text, fill, color) in enumerate(boxes):
        x = left0 + i * (bw + gap)
        _box(slide, Emu(x), Emu(top), Emu(bw), Emu(bh), text, fill, color, size=13, bold=(i == 3))
        if i < n - 1:
            _arrow(slide, Emu(int(x + bw + gap*0.15)), Emu(int(top + bh/2)), Emu(int(x + bw + gap*0.85)), Emu(int(top + bh/2)))

    _caption(slide, MARGIN, Emu(4300000), Emu(10800000),
             "Un score de moins l'infini ne peut jamais etre le maximum : le coup illegal ne peut donc jamais etre choisi, quelle que soit la position.",
             size=13, italic=False)
    _footer(slide, "Architecture")
    _note(slide, "Choisir un coup legal", (
        "Question que le jury pose presque toujours : est-ce que le bot peut proposer un coup "
        "impossible ? La reponse est non, et c'est garanti par construction, pas par une verification "
        "a posteriori. Le reseau calcule un score pour les 1968 coups du vocabulaire, y compris ceux "
        "qui n'ont aucun sens dans la position actuelle. Avant de choisir, on repere tous les coups "
        "illegaux dans cette position precise, et on remplace leur score par moins l'infini. Un score "
        "de moins l'infini ne peut mathematiquement jamais etre le plus grand : le coup illegal est "
        "donc ecarte de facon absolue, quelle que soit la position, sans exception possible."
    ))
    return slide


# ---------------------------------------------------------------------------
# Contenu des diapos texte / tableau
# ---------------------------------------------------------------------------

def build(prs):
    slide_cover(
        prs, "Chess Bot v1", "Un Transformer qui joue aux echecs sans recherche",
        ["Projet de substitution au stage", "SEIBOU Naadjath & LAKRA Rajaa", "ECE, soutenance semaine du 31 aout 2026"],
        "Bonjour. Nous presentons Chess Bot v1 : un reseau de neurones qui choisit un "
        "coup d'echecs directement a partir de la position, sans explorer la moindre "
        "variante future. C'est cette absence de recherche qui rend le projet interessant, "
        "et on va vous montrer ce qu'un petit modele arrive a apprendre, et surtout ses "
        "limites. Je suis Naadjath, voici Rajaa, on se repartit la presentation."
    )

    slide_bullets(
        prs, "La question de depart", "", kicker="Introduction",
        bullets=[
            ("Les moteurs classiques CHERCHENT.",
             "Deep Blue, Stockfish : ils explorent des millions de positions futures avant de jouer un coup."),
            ("Notre bot, lui, ne cherche jamais.",
             "Il repond a l'instinct, comme un joueur humain fort en partie a la pendule tres rapide (blitz)."),
            ("La question qu'on se pose :",
             "un Transformer peut-il jouer correctement rien qu'en RECONNAISSANT une position, sans jamais calculer une suite de coups ?"),
        ],
        oral=(
            "Pour comprendre l'interet du projet, il faut d'abord voir ce qui rend les moteurs "
            "d'echecs actuels si forts. Deep Blue, Stockfish, tous cherchent : a chaque coup, ils "
            "explorent mentalement des millions de positions futures avant de se decider. C'est de la "
            "force brute, tres efficace, mais couteuse en calcul. Nous posons une question differente. "
            "Est-ce qu'un reseau de neurones peut jouer correctement sans jamais faire ce travail de "
            "recherche, juste en reconnaissant une position comme un joueur humain reconnait un motif "
            "familier ? C'est cette question qui guide tout le reste de la presentation."
        ),
        tag="Introduction",
    )

    slide_bullets(
        prs, "Contexte : les echecs et l'IA", "Un demi-siecle d'approches, un point commun : la recherche", kicker="Etat de l'art",
        bullets=[
            ("1997 : Deep Blue bat Kasparov", "recherche alpha-beta massive, evaluation ecrite a la main par des experts"),
            ("2017 : AlphaZero", "reseau de neurones + recherche arborescente (MCTS), entraine par auto-apprentissage"),
            ("2020+ : Stockfish (NNUE)", "recherche alpha-beta tres optimisee + petit reseau d'evaluation"),
            ("2024 : DeepMind, Grandmaster-Level Chess Without Search", "un Transformer SEUL atteint un niveau de grand maitre, c'est notre reference directe"),
        ],
        oral=(
            "Petit rappel historique pour situer le projet. En 1997, Deep Blue bat Kasparov grace a "
            "une recherche massive et une evaluation ecrite a la main par des experts. En 2017, "
            "AlphaZero associe un reseau de neurones a une recherche arborescente, et apprend en "
            "jouant contre lui-meme des millions de fois. Aujourd'hui, Stockfish, le moteur le plus "
            "fort au monde, combine toujours une recherche tres optimisee et un petit reseau "
            "d'evaluation. A chaque fois, le point commun, c'est la recherche. Ce qui change tout en "
            "2024, c'est cet article de DeepMind qui montre qu'un Transformer seul, sans la moindre "
            "recherche, peut atteindre un niveau de grand maitre. C'est notre reference directe : on "
            "s'inspire de cette demarche, a une echelle beaucoup plus modeste."
        ),
        tag="Etat de l'art",
    )

    slide_bullets(
        prs, "Notre pari : une version reduite et assumee", "", kicker="Positionnement",
        bullets=[
            ("DeepMind (2024)", "270 millions de parametres · 10 millions de parties annotees par Stockfish · TPU dedies"),
            ("Nous", "6,86 millions de parametres (x40 plus petit) · 1 million de positions · 1 GPU T4 gratuit (Google Colab)"),
            ("Consequence assumee",
             "on ne vise pas a battre Stockfish. L'objectif est de mesurer honnetement ce qu'un petit modele apprend, avec la meme rigueur experimentale."),
        ],
        oral=(
            "Il faut etre honnete sur l'echelle du projet des le depart. DeepMind travaille avec 270 "
            "millions de parametres, 10 millions de parties annotees par Stockfish, et des processeurs "
            "specialises TPU. Nous avons un modele 40 fois plus petit, 1 million de positions au lieu "
            "de 10, et un unique GPU gratuit accessible depuis un navigateur. La consequence, on "
            "l'assume completement : on ne cherche pas a battre Stockfish, et personne ne s'attend a "
            "ce qu'on y arrive. Notre objectif est different, et plus modeste : mesurer honnetement ce "
            "qu'un petit modele, avec des moyens accessibles a des etudiants, arrive reellement a "
            "apprendre, avec la meme rigueur scientifique qu'un vrai laboratoire."
        ),
        tag="Positionnement",
    )

    diagram_pipeline(prs)
    diagram_funnel(prs)
    diagram_encoding(prs)
    diagram_architecture(prs)
    diagram_masking(prs)

    slide_image(
        prs, "L'entrainement", "4 epoques sur 1 million de positions, GPU T4 gratuit (Google Colab)",
        ASSETS / "courbes_entrainement.png",
        "La perte descend regulierement sous le niveau du hasard (ligne pointillee), sans surapprentissage : "
        "les courbes d'entrainement et de validation restent proches.",
        oral=(
            "Voici les courbes reelles de notre entrainement, pas des chiffres illustratifs. A gauche, "
            "la perte : elle part de 7,6, qui correspond au hasard pur sur 1968 coups possibles, c'est "
            "le logarithme de 1968, et descend a 3,08 apres 4 epoques. Au milieu, la top-1 : la "
            "capacite du modele a retrouver exactement le coup joue par l'humain. Elle atteint 22,7 % "
            "en fin d'entrainement. A droite, la top-5 : le bon coup est dans les 5 premieres "
            "propositions du modele 51,7 % du temps. Point important a expliquer au jury si la "
            "question vient : 22 % n'est pas un mauvais score, car dans beaucoup de positions plusieurs "
            "coups se valent. Cette mesure evalue l'imitation, pas la force de jeu reelle, c'est pour "
            "ca qu'on mesure aussi l'Elo par de vraies parties, sur la diapo suivante. Chaque epoque "
            "prend 8 a 9 minutes sur le GPU gratuit."
        ),
        kicker="Entrainement", tag="Resultats",
    )

    slide_image(
        prs, "Il a appris les ouvertures", "Sans qu'aucune regle ne lui ait ete enseignee",
        ASSETS / "ouverture_echiquier.png",
        "Dans la position de depart, le reseau concentre ses propositions sur le developpement des "
        "cavaliers et l'occupation du centre : des principes d'ouverture reels.",
        oral=(
            "C'est le resultat le plus parlant du projet, celui qu'on garde pour la fin de la partie "
            "technique. Si on demande au reseau ce qu'il joue au tout premier coup, avant meme d'avoir "
            "vu un seul coup de la partie, il repond d3, developper le cavalier en c3 ou en f3, jouer "
            "d4 ou c4. Ce sont exactement les principes qu'on enseigne a un debutant : sortir ses "
            "pieces, prendre le centre. Et voici le point cle : le modele n'a jamais recu une seule "
            "regle des echecs, ni comment une piece se deplace, ni ce qu'est un echec et mat. Il a "
            "uniquement regarde des parties de joueurs forts et en a deduit ces principes tout seul. "
            "C'est la preuve concrete qu'il a appris quelque chose de reel, et pas seulement memorise "
            "des positions par coeur."
        ),
        kicker="Resultats", tag="Resultats",
    )

    slide_image(
        prs, "Le niveau mesure en parties reelles", "60 parties par adversaire, couleurs alternees, intervalles de confiance de Wilson (95 %)",
        ASSETS / "elo_resultats.png",
        "Face a Stockfish bride a son niveau minimum (1320), le bot perd toutes ses parties : "
        "son niveau reel est nettement en dessous.",
        oral=(
            "Deuxieme resultat, et le plus important pour juger honnetement le projet : le niveau reel, "
            "mesure en faisant vraiment jouer le bot, pas en extrapolant depuis la top-1. On l'a fait "
            "affronter une echelle d'adversaires de force croissante, 60 parties chacun, avec un "
            "intervalle de confiance calcule par la methode de Wilson : on ne donne jamais un chiffre "
            "seul, toujours une fourchette. Contre l'adversaire aleatoire, le bot obtient 47,5 %, un "
            "score statistiquement equivalent au hasard, autour de 230 Elo. Contre tous les adversaires "
            "plus structures, il perd presque systematiquement. Face a Stockfish bride a 1320, le "
            "reglage le plus faible que Stockfish accepte, il ne gagne aucune partie sur les 60 jouees. "
            "Notre verdict honnete : le bot se situe autour de 250 a 300 Elo, un niveau de tout premier "
            "debutant. On prefere assumer ce resultat plutot que de le maquiller."
        ),
        kicker="Resultats", tag="Resultats", max_h=Emu(4300000),
    )

    slide_bullets(
        prs, "Analyse : ce que ca revele", "", kicker="Discussion",
        bullets=[
            ("Force : la connaissance des principes",
             "developpement, occupation du centre, acquis par simple observation, sans recherche."),
            ("Faiblesse : un jeu passif",
             "51 nulles sur 60 parties contre l'aleatoire. Le bot atteint des positions gagnantes mais ne les convertit pas, il tourne en rond."),
            ("Interpretation",
             "sans recherche, le modele ne verifie jamais une suite de coups. Il reconnait des schemas, mais ne calcule pas les consequences. Un entrainement plus long (les courbes montaient encore a la 4e epoque) aurait probablement ameliore ce point."),
        ],
        oral=(
            "Que retenir de ces deux resultats mis cote a cote ? Il y a une vraie force : la "
            "connaissance des principes de base, developpement des pieces, occupation du centre, "
            "acquise sans la moindre regle explicite. Mais il y a aussi une faiblesse nette : un jeu "
            "passif. Cinquante et une nulles sur soixante parties contre l'aleatoire, c'est enorme. Le "
            "bot atteint souvent des positions gagnantes, mais il ne sait pas les convertir, il tourne "
            "en rond jusqu'a la regle des 50 coups ou la triple repetition. Notre interpretation : sans "
            "recherche, le modele reconnait des schemas mais ne verifie jamais une suite de coups a "
            "l'avance. Il sait QUOI faire au debut, pas COMMENT finir. Et les courbes montraient "
            "encore une progression a la quatrieme epoque, donc un entrainement plus long aurait "
            "probablement attenue ce probleme."
        ),
        tag="Discussion",
    )

    slide_bullets(
        prs, "L'application", "Jouer contre le bot, et voir ce qu'il pense", kicker="Livrable",
        bullets=[
            ("Mode Jouer", "on affronte le Transformer (ou les adversaires de reference) dans le navigateur."),
            ("Mode spectateur", "deux bots s'affrontent automatiquement, on regarde et on commente."),
            ("Le detail cle", "a chaque coup, les probabilites reellement sorties du reseau s'affichent : on voit litteralement le modele reflechir."),
        ],
        oral=(
            "On a aussi construit une application de jeu complete, dans le navigateur, sans aucune "
            "dependance a installer. Deux facons de l'utiliser : le mode Jouer, ou on affronte "
            "directement le Transformer ou n'importe quel bot de reference, et le mode spectateur, ou "
            "deux bots s'affrontent tout seuls pendant qu'on regarde et qu'on commente, pratique pour "
            "une demonstration sans risquer de se tromper de coup en direct. Le detail le plus "
            "interessant : a chaque coup joue, les probabilites reellement calculees par le reseau "
            "s'affichent a l'ecran. On ne voit pas juste le resultat, on voit litteralement le modele "
            "hesiter entre plusieurs coups en temps reel. [Montrer une courte demonstration ici si "
            "possible.]"
        ),
        tag="Demonstration",
    )

    slide_bullets(
        prs, "Limites et perspectives", "", kicker="Discussion",
        bullets=[
            ("Passer a l'action-value", "predire la probabilite de gagner de chaque coup, comme DeepMind, plutot que d'imiter le coup humain."),
            ("Annoter avec Stockfish", "apprendre du MEILLEUR coup calcule par un moteur, pas seulement du coup joue par un humain."),
            ("Ajouter une recherche legere", "une profondeur 2-3 guidee par les propositions du reseau corrigerait probablement l'essentiel des erreurs tactiques."),
            ("Entrainer plus longtemps", "les courbes etaient encore croissantes ; un GPU stable (non sujet aux deconnexions) permettrait d'aller plus loin."),
        ],
        oral=(
            "Avec plus de temps, quatre pistes concretes permettraient d'ameliorer le bot. D'abord, "
            "changer d'objectif d'entrainement : predire la probabilite de gagner de chaque coup, "
            "comme fait DeepMind, plutot que d'imiter simplement le coup joue par l'humain. Ensuite, "
            "faire annoter nos donnees par Stockfish, pour apprendre du meilleur coup possible et pas "
            "seulement de celui qu'un humain a choisi, qui peut etre imparfait. Troisieme piste : "
            "ajouter une petite recherche, profondeur 2 ou 3, guidee par les propositions du reseau, ce "
            "qui corrigerait probablement l'essentiel des erreurs tactiques que le bot commet "
            "actuellement. Et enfin, tout simplement entrainer plus longtemps : nos courbes montaient "
            "encore a la fin, la limite venait des deconnexions de l'environnement gratuit, pas d'un "
            "plafond du modele."
        ),
        tag="Perspectives",
    )

    slide_bullets(
        prs, "Gestion de projet", "Un binome, un depot Git versionne", kicker="Organisation",
        bullets=[
            ("Naadjath & Rajaa", "travail mixte sur l'ensemble de la chaine plutot qu'une repartition stricte par axe : donnees, modele, evaluation et application ont ete avancees en binome."),
            ("Suivi", "Git/GitHub avec des commits reguliers des deux membres, ce qui retrace la progression reelle du projet."),
            ("Difficultes reelles rencontrees",
             "lecture des donnees initialement trop lente (corrigee, x18) · Elo infini a 100 % de victoires (corrige par la methode de Wilson) · deconnexions de l'environnement Colab gratuit (contournees par sauvegarde sur Drive a chaque epoque)."),
        ],
        oral=(
            "Un mot sur l'organisation du binome. On n'a pas fige un partage strict des taches : on a "
            "avance en binome sur l'ensemble de la chaine, donnees, modele, evaluation et application, "
            "en se repartissant le travail au fur et a mesure selon les priorites. Le tout est versionne "
            "sur Git, avec un historique de commits des deux membres qui retrace la progression. On a "
            "rencontre de vraies difficultes techniques, pas seulement conceptuelles : la lecture des "
            "donnees etait au depart 18 fois trop lente, corrigee en changeant de strategie de filtrage ; "
            "le calcul d'Elo donnait un resultat infini a 100 % de victoires, corrige avec la methode de "
            "Wilson ; et l'environnement Colab gratuit se deconnectait regulierement, ce qu'on a "
            "contourne en sauvegardant le modele sur Google Drive a chaque epoque."
        ),
        tag="Organisation",
    )

    slide_bullets(
        prs, "Conclusion", "", kicker="Synthese",
        bullets=[
            ("Un Transformer apprend a jouer sans recherche", "juste en observant des parties, par clonage comportemental."),
            ("Niveau modeste (~250 Elo) mais reel", "principes d'ouverture correctement acquis, jeu de fin de partie encore faible."),
            ("Une demarche complete et rigoureuse", "chaine reproductible, encodage teste, evaluation avec intervalles de confiance."),
            ("Le point cle a retenir", "la difference entre connaissance (ce que le reseau a memorise) et calcul (la recherche qu'on a volontairement retiree)."),
        ],
        oral=(
            "Pour conclure. Nous avons montre qu'un Transformer, meme modeste, apprend a jouer aux "
            "echecs sans jamais explorer une seule variante future, juste en observant des parties "
            "d'un bon niveau. Le niveau atteint reste faible, autour de 250 Elo, mais il est reel : le "
            "modele a correctement acquis des principes d'ouverture, meme si sa fin de partie reste "
            "fragile. Surtout, la demarche est complete et rigoureuse du debut a la fin : chaine de "
            "traitement reproductible, encodage verifie par des tests automatiques, evaluation "
            "toujours accompagnee d'un intervalle de confiance. Le point que nous voulons que vous "
            "reteniez : ce projet illustre concretement la difference entre la connaissance, ce que le "
            "reseau a memorise en observant, et le calcul, la recherche qu'on a volontairement "
            "retiree. C'est cette absence de calcul qui explique a la fois ce que le modele reussit et "
            "ce qu'il rate."
        ),
        tag="Conclusion",
    )

    slide_cover(
        prs, "Merci", "Questions ?", [],
        "Merci de votre attention, nous sommes prets pour vos questions. (Reponses preparees : "
        "pourquoi un Transformer plutot qu'un reseau convolutif : l'attention relie directement deux "
        "cases eloignees, utile pour une piece a longue portee. Comment on garantit qu'aucun coup "
        "illegal n'est joue : masquage a moins l'infini avant le choix du maximum. Pourquoi l'Elo "
        "varie autant selon l'adversaire : les Elo des baselines sont des ordres de grandeur admis, "
        "pas des valeurs calibrees officiellement.)"
    )


def build_pptx() -> None:
    _ORAL_SCRIPT.clear()  # une nouvelle construction ne doit jamais accumuler sur la precedente
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    build(prs)
    n_slides = len(prs.slides.__iter__.__self__._sldIdLst)
    if len(_ORAL_SCRIPT) != n_slides:
        raise AssertionError(
            f"{n_slides} diapos mais {len(_ORAL_SCRIPT)} textes oraux enregistres : "
            "une diapo a ete ajoutee sans son commentaire oral (verifier les derniers appels a _note)."
        )
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_OUT)
    print(f"Diaporama ecrit : {PPTX_OUT}  ({n_slides} diapos)")


def build_oral_guide() -> None:
    """Ecrit le guide oral a partir du texte enregistre pendant build(),
    jamais en relisant le pptx, pour ne perdre ni un titre ni un commentaire.
    """
    lines = [
        "# Guide oral : soutenance Chess Bot v1",
        "",
        "*Le script complet, diapo par diapo. A imprimer et garder en main.*",
        "",
        "**Conseils :**",
        "- Ne lisez pas la diapo au jury : elle est pour lui, ce texte est pour vous.",
        "- Repetez a voix haute au moins 3 fois en chronometrant (visez 12-15 minutes).",
        "- Alternez les diapos entre les deux membres du binome.",
        "- Les passages entre crochets, comme [Prenom] ou [X], sont a personnaliser ou a jouer en direct.",
        "",
        "---",
        "",
    ]
    for i, (title, oral) in enumerate(_ORAL_SCRIPT, 1):
        lines.append(f"## Diapo {i} : {title}")
        lines.append("")
        lines.append(f"> {oral}")
        lines.append("")
        lines.append("---")
        lines.append("")

    ORAL_OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Guide oral ecrit : {ORAL_OUT}  ({len(_ORAL_SCRIPT)} diapos commentees)")


if __name__ == "__main__":
    if not ASSETS.exists() or not any(ASSETS.iterdir()):
        raise SystemExit("Lancez d'abord : python tools/make_slide_assets.py")
    build_pptx()
    build_oral_guide()
